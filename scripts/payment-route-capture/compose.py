"""
캡처된 PNG 들을 토스 결제경로 가이드 순서로 PPT (16:9) 슬라이드 조립.

표지 데이터는 cover_info.json (없으면 cover_info.example.json) 사용.
.env.local 의 NEXT_PUBLIC_BUSINESS_COMPANY_NAME / REG_NUMBER 가 있으면 우선.
"""

from __future__ import annotations

import json
import os
import re
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from config import OUTPUT_DIR, PPTX_PATH, SCREENSHOT_DIR

ROOT = Path(__file__).resolve().parent
ENV_FILE = ROOT.parent.parent / "next" / ".env.local"
COVER_JSON = ROOT / "cover_info.json"
COVER_EXAMPLE = ROOT / "cover_info.example.json"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 8.889  # 1920x1280 (3:2) 캡처 fit

# 슬라이드 순서 — (캡처 키, 슬라이드 헤더)
SLIDE_PLAN: list[tuple[str, str]] = [
    ("02_footer",              "② 하단 정보 — 사업자 6종 표시"),
    ("03_refund",              "③-1 환불 정책 — 취소·반품·교환 기준"),
    ("03_refund_2",            "③-2 환불 처리 — 배송비 부담 + 환불 절차"),
    ("04a_login",              "④-A 로그인"),
    ("04b_signup",             "④-B 회원가입"),
    ("04c_guest_checkout",     "④-C 비회원 구매 경로"),
    ("05a_shop",               "⑤-1 상품 목록 (/shop)"),
    ("05b_pdp",                "⑤-2 상품 상세"),
    ("05c_cart",               "⑤-3 장바구니"),
    ("05d_checkout_member",    "⑤-4 주문서 (회원)"),
    ("05e_checkout_guest",     "⑤-4 주문서 (비회원)"),
    ("06a_widget_card",        "⑥-1 결제수단 — 카드 선택"),
    ("06b_widget_card_select", "⑥-2 카드사 드롭다운"),
    ("06c_bc_auth",            "⑥-3 비씨카드 인증 화면"),
]


@dataclass
class CoverInfo:
    company_name: str
    registration_number: str
    url: str
    test_id: str
    test_pw: str
    note: str


def load_env_business() -> dict[str, str]:
    """next/.env.local 에서 NEXT_PUBLIC_BUSINESS_* 추출 (값에 # 주석은 제외)."""
    if not ENV_FILE.exists():
        return {}
    pattern = re.compile(r"^(NEXT_PUBLIC_BUSINESS_[A-Z_]+)=(.*)$")
    out: dict[str, str] = {}
    for line in ENV_FILE.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        m = pattern.match(line)
        if m:
            value = m.group(2).strip().strip('"').strip("'")
            out[m.group(1)] = value
    return out


def load_cover() -> CoverInfo:
    src = COVER_JSON if COVER_JSON.exists() else COVER_EXAMPLE
    data = json.loads(src.read_text(encoding="utf-8"))
    env = load_env_business()
    return CoverInfo(
        company_name=env.get("NEXT_PUBLIC_BUSINESS_COMPANY_NAME", data.get("company_name", "")),
        registration_number=env.get("NEXT_PUBLIC_BUSINESS_REG_NUMBER", data.get("registration_number", "")),
        url=data.get("url", "https://goodthings-roasters.com"),
        test_id=data.get("test_id", ""),
        test_pw=data.get("test_pw", ""),
        note=data.get("note", ""),
    )


def setup_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)


def add_cover_slide(prs: Presentation, info: CoverInfo) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # 타이틀
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11.3), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "굳띵즈 로스터스 — 결제경로 제작 가이드"
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True

    # 정보 박스
    rows = [
        ("(1) 상호명",      info.company_name),
        ("(2) 사업자번호",  info.registration_number),
        ("(3) URL",         info.url),
        ("(4) Test ID",     info.test_id),
        ("(5) Test PW",     info.test_pw),
    ]
    if info.note:
        rows.append(("비고", info.note))

    y = 2.3
    for label, value in rows:
        lb = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(3.0), Inches(0.55))
        lb.text_frame.text = label
        for r in lb.text_frame.paragraphs[0].runs:
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)  # 토스 PDF 톤
        vb = slide.shapes.add_textbox(Inches(5.2), Inches(y), Inches(7.0), Inches(0.55))
        vb.text_frame.text = value or "—"
        for r in vb.text_frame.paragraphs[0].runs:
            r.font.size = Pt(20)
        y += 0.65


def add_image_slide(prs: Presentation, image_path: Path, title: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # 제목 (상단)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].runs[0].font.size = Pt(20)
    tf.paragraphs[0].runs[0].font.bold = True

    # 이미지 — 1920x1280 (3:2) 캡처에 맞춰 슬라이드 안 fit.
    # box 12.0 x 8.0 (3:2 비율 유지) + top 0.85 = bottom 8.85 < slide 8.889 → 넘침 없음.
    img_box_w = Inches(12.0)
    img_box_h = Inches(8.0)  # 12.0 / 1.5 (3:2 비율 유지)
    left = (prs.slide_width - img_box_w) / 2
    top = Inches(0.85)
    slide.shapes.add_picture(str(image_path), left, top, width=img_box_w, height=img_box_h)


def add_thanks_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0), Inches(3.0), Inches(SLIDE_W_IN), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "감사합니다"
    p.alignment = 2  # center
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.bold = True


def compose() -> None:
    info = load_cover()
    prs = Presentation()
    setup_slide_size(prs)

    add_cover_slide(prs, info)

    missing: list[str] = []
    for key, title in SLIDE_PLAN:
        img = SCREENSHOT_DIR / f"{key}.png"
        if img.exists():
            add_image_slide(prs, img, title)
        else:
            missing.append(key)

    add_thanks_slide(prs)

    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)

    print(f"\nPPT 저장: {PPTX_PATH}")
    print(f"슬라이드 총 {len(prs.slides)}개 (표지 1 + 본문 {len(SLIDE_PLAN) - len(missing)} + 감사 1)")
    if missing:
        print("\n[누락된 캡처 — 수동 추가 또는 capture.py 재실행 권장]")
        for k in missing:
            print(f"  - {k}.png")


if __name__ == "__main__":
    compose()
