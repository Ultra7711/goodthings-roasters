"""
정기결제(빌링) PPT 조립 — compose.py 답습.

차이점:
- 표지 타이틀: "결제경로 제작 가이드 (빌링)"
- SLIDE_PLAN: 본문 A (일반결제 재사용 5) + B (정기결제 신규 7) + C (어필 3) = 15장
- 출력: output/payment-route-billing.pptx

원본 compose.py 의 헬퍼 (load_env_business / load_cover / setup_slide_size / add_image_slide / add_thanks_slide) 재사용.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from compose import (
    SLIDE_H_IN,
    SLIDE_W_IN,
    CoverInfo,
    add_image_slide,
    add_thanks_slide,
    load_cover,
    setup_slide_size,
)
from config import OUTPUT_DIR, SCREENSHOT_DIR

BILLING_PPTX_PATH = OUTPUT_DIR / "payment-route-billing.pptx"

# 슬라이드 순서 — (캡처 키, 슬라이드 헤더)
SLIDE_PLAN: list[tuple[str, str]] = [
    # 본문 A — 일반결제 재사용 (5)
    ("02_footer",                       "② 하단 정보 — 사업자 6종 표시"),
    ("04a_login",                       "④-A 로그인"),
    ("04b_signup",                      "④-B 회원가입"),
    ("04c_guest_checkout",              "④-C 비회원 구매 경로"),
    ("05a_shop",                        "⑤-1 상품 목록 (/shop)"),
    # 본문 B — 정기결제 신규 (7)
    ("b_03_refund_returns",             "③-1 환불 정책 — 정기배송 (/legal/returns)"),
    ("b_03_refund_faq",                 "③-2 정기결제 Q&A (/legal/payment-faq)"),
    ("b_05b_pdp_subscription",          "⑤-2 상품 상세 — 정기배송 토글 ON"),
    ("b_05c_cart_subscription",         "⑤-3 장바구니 — 정기배송 라벨"),
    ("b_05d_checkout_subscription",     "⑤-4 주문서 — 정기결제 분기 (γ 합산 · mixed cart)"),
    ("b_06a_billing_card",              "⑥-1 빌링 카드 정보 입력창 (CARD)"),
    ("b_06b_billing_transfer",          "⑥-2 빌링 계좌이체 입력창 (TRANSFER)"),
    # 본문 C — 어필 (3)
    ("b_app_a1_pdp_normal",             "어필-1 PDP — 일반결제 상품 옵션"),
    ("b_app_a2_pdp_subscription",       "어필-2 PDP — 정기배송 토글 ON (분기 비교)"),
    ("b_app_b_terms_no_lockin",         "어필-3 TERMS §10조의2 — 약정 없음 + 즉시 해지"),
]


def add_billing_cover_slide(prs: Presentation, info: CoverInfo) -> None:
    """표지 — '결제경로 제작 가이드 (빌링)' 타이틀 + 가맹점 정보."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # 타이틀 — 일반결제와 구분
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11.3), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "굳띵즈 로스터스 — 결제경로 제작 가이드 (빌링)"
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True

    # 정보 박스
    rows = [
        ("(1) 상호명",      info.company_name),
        ("(2) 사업자번호",  info.registration_number),
        ("(3) URL",         info.url),
        ("(4) Test ID",     info.test_id),
        ("(5) Test PW",     info.test_pw),
    ]
    if info.note:
        rows.append(("비고", info.note))

    y = 2.3
    for label, value in rows:
        lb = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(3.0), Inches(0.55))
        lb.text_frame.text = label
        for r in lb.text_frame.paragraphs[0].runs:
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)  # 토스 PDF 톤
        vb = slide.shapes.add_textbox(Inches(5.2), Inches(y), Inches(7.0), Inches(0.55))
        vb.text_frame.text = value or "—"
        for r in vb.text_frame.paragraphs[0].runs:
            r.font.size = Pt(20)
        y += 0.65


def compose() -> None:
    info = load_cover()
    prs = Presentation()
    setup_slide_size(prs)

    add_billing_cover_slide(prs, info)

    missing: list[str] = []
    for key, title in SLIDE_PLAN:
        img = SCREENSHOT_DIR / f"{key}.png"
        if img.exists():
            add_image_slide(prs, img, title)
        else:
            missing.append(key)

    add_thanks_slide(prs)

    BILLING_PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(BILLING_PPTX_PATH)

    print(f"\n정기결제 PPT 저장: {BILLING_PPTX_PATH}")
    print(f"슬라이드 총 {len(prs.slides)}개 (표지 1 + 본문 {len(SLIDE_PLAN) - len(missing)} + 감사 1)")
    if missing:
        print("\n[누락된 캡처 — images/screenshots/ 에 다음 파일 추가 필요]")
        for k in missing:
            print(f"  - {k}.png")


if __name__ == "__main__":
    compose()
